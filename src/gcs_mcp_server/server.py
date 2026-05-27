import os
import argparse
from mcp.server.fastmcp import FastMCP
from google.cloud import storage
from google.oauth2 import service_account

mcp = FastMCP("gcs-mcp-server")

# グローバルクライアント（main()で初期化される）
client: storage.Client | None = None


def build_client(project: str, key_file: str | None) -> storage.Client:
    if key_file:
        credentials = service_account.Credentials.from_service_account_file(key_file)
        return storage.Client(project=project, credentials=credentials)
    return storage.Client(project=project)


# ===================================================================
# 読み取り系
# ===================================================================

@mcp.tool()
def list_buckets() -> list[str]:
    """プロジェクト内のバケット一覧を返す"""
    return [b.name for b in client.list_buckets()]


@mcp.tool()
def list_files(bucket: str, prefix: str = "") -> list[dict]:
    """
    バケット内のファイル一覧を返す

    Args:
        bucket: バケット名
        prefix: パスのプレフィックス（例: "data/2024/"）省略すると全件
    """
    blobs = client.list_blobs(bucket, prefix=prefix)
    return [
        {
            "name":    blob.name,
            "size":    blob.size,
            "updated": blob.updated.isoformat() if blob.updated else None,
        }
        for blob in blobs
    ]


@mcp.tool()
def read_file(bucket: str, name: str, max_bytes: int = 100 * 1024 * 1024) -> str:
    """
    ファイルの内容をテキストとして読み込む

    Args:
        bucket:    バケット名
        name:      ファイルパス（例: "data/report.csv"）
        max_bytes: 読み込みを許可する最大バイト数（デフォルト: 100MB）
    """
    blob = client.bucket(bucket).blob(name)
    blob.reload()
    if blob.size > max_bytes:
        raise ValueError(
            f"ファイルが大きすぎます ({blob.size:,} bytes > {max_bytes:,} bytes)"
        )
    return blob.download_as_text(encoding="utf-8")


@mcp.tool()
def get_file_info(bucket: str, name: str) -> dict:
    """
    ファイルのメタデータを返す（サイズ・更新日時・Content-Typeなど）

    Args:
        bucket: バケット名
        name:   ファイルパス
    """
    blob = client.bucket(bucket).blob(name)
    blob.reload()
    return {
        "name":         blob.name,
        "size":         blob.size,
        "content_type": blob.content_type,
        "updated":      blob.updated.isoformat() if blob.updated else None,
        "md5_hash":     blob.md5_hash,
    }


# ===================================================================
# 書き込み系
# ===================================================================

@mcp.tool()
def read_pptx(bucket: str, name: str) -> str:
    """
    GCSのPowerPointファイル（.pptx）を読み込み、スライドのテキストを返す

    Args:
        bucket: バケット名
        name:   ファイルパス（例: "data/presentation.pptx"）
    """
    import io
    from pptx import Presentation

    data = client.bucket(bucket).blob(name).download_as_bytes()
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            slides.append(f"## スライド {i}\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "（テキストなし）"


@mcp.tool()
def download_file(bucket: str, name: str, local_path: str) -> str:
    """
    GCSのファイルをローカルパスに保存する

    Args:
        bucket:     バケット名
        name:       ダウンロード元ファイルパス（例: "data/report.csv"）
        local_path: 保存先のローカルファイルパス（例: "/tmp/report.csv"）
    """
    import pathlib
    resolved = pathlib.Path(local_path).expanduser().resolve()
    resolved.parent.mkdir(parents=True, exist_ok=True)
    client.bucket(bucket).blob(name).download_to_filename(str(resolved))
    size = resolved.stat().st_size
    return f"ダウンロード完了: gs://{bucket}/{name} → {resolved} ({size:,} bytes)"


@mcp.tool()
def upload_text(bucket: str, name: str, content: str, content_type: str = "text/plain") -> str:
    """
    テキストをファイルとしてアップロードする

    Args:
        bucket:       バケット名
        name:         保存先ファイルパス（例: "output/result.txt"）
        content:      ファイルの内容（テキスト）
        content_type: MIMEタイプ（デフォルト: text/plain）
    """
    blob = client.bucket(bucket).blob(name)
    blob.upload_from_string(content, content_type=content_type)
    return f"アップロード完了: gs://{bucket}/{name}"


@mcp.tool()
def copy_file(src_bucket: str, src_name: str, dst_bucket: str, dst_name: str) -> str:
    """
    ファイルを別の場所にコピーする

    Args:
        src_bucket: コピー元バケット名
        src_name:   コピー元ファイルパス
        dst_bucket: コピー先バケット名
        dst_name:   コピー先ファイルパス
    """
    src_blob = client.bucket(src_bucket).blob(src_name)
    client.copy_blob(src_blob, client.bucket(dst_bucket), dst_name)
    return f"コピー完了: gs://{src_bucket}/{src_name} → gs://{dst_bucket}/{dst_name}"


@mcp.tool()
def delete_file(bucket: str, name: str) -> str:
    """
    ファイルを削除する

    Args:
        bucket: バケット名
        name:   削除するファイルパス
    """
    client.bucket(bucket).blob(name).delete()
    return f"削除完了: gs://{bucket}/{name}"


# ===================================================================
# エントリーポイント
# ===================================================================

def main():
    parser = argparse.ArgumentParser(description="GCS MCP Server")
    parser.add_argument("--project",  required=True, help="GCPプロジェクトID")
    parser.add_argument("--key-file", default=None,  help="サービスアカウントJSONのパス（省略時はADC）")
    args = parser.parse_args()

    global client
    client = build_client(args.project, args.key_file)

    mcp.run(transport="stdio")


if __name__ == "__main__":
    main()
